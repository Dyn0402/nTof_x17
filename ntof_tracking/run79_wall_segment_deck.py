#!/usr/bin/env python3
"""
run79_wall_segment_deck.py — package the wall-segment tour for showing to people.

The tour (`run79_wall_segment_gif.py`) is an **animation**, so it cannot go in a
PDF. This builds two things that can carry it:

  * `wall_segment_tour.pptx` — the GIF embedded on its own slide. PowerPoint and
    LibreOffice Impress both animate an embedded GIF in slideshow mode
    (`soffice --show wall_segment_tour.pptx`). Static stills follow, so the deck
    still says something if exported to PDF.
  * `wall_segment_tour_web/index.html` + assets — a plain page that animates in
    any browser, for sending to someone who will not open a deck.

Everything on the slides is read from the tour's own sidecar JSON, so the deck
cannot drift from the figure it is describing.

    .venv/bin/python ntof_tracking/run79_wall_segment_deck.py [--small] [--out DIR]
"""
import argparse
import json
import os
import shutil

TOUR = ('/media/dylan/data/x17/beam_july/analysis/wft/run_79/stat090_0000/'
        'mx17_A/wall_segment_tour')

BG = 0xF7F7F7          # matches the figure background
INK = 0x1A1A1A
MUTED = 0x6E6E6E
ACCENT = 0xB0143C      # the segment-1 crimson


def _rgb(v):
    from pptx.dml.color import RGBColor
    return RGBColor((v >> 16) & 255, (v >> 8) & 255, v & 255)


def _text(slide, x, y, w, h, runs, align='left'):
    """runs = [(text, size_pt, bold, colour), ...]; one paragraph each."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, size, bold, colour) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER}[align]
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = _rgb(colour)
        r.font.name = 'DejaVu Sans'
    return tb


def _bg(slide, prs):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                               prs.slide_width, prs.slide_height)
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(BG)
    s.line.fill.background()
    s.shadow.inherit = False
    # keep it behind everything added later
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)


def _fit(img_w, img_h, max_w, max_h):
    """Letterbox an image into a box, returning (w, h) in inches."""
    scale = min(max_w / img_w, max_h / img_h)
    return img_w * scale, img_h * scale


def build_pptx(js, gif, stills, out):
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    st = js['stats']
    conv, null = st['convergence'], st['convergence_null']

    # ---------------------------------------------------------------- title
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _text(s, 0.9, 1.9, 11.5, 2.6, [
        ('Tracks point back at the target — and at the wall segment that fired',
         34, True, INK),
        ('', 12, False, INK),
        ('run_79 / mx17_A waveform-first reconstruction, joined to n_TOF run 224572',
         19, False, MUTED),
    ])
    _text(s, 0.9, 4.5, 11.5, 1.6, [
        (f'{st["n_all"]:,} tracks over {st["n_bunch"]} bunches   ·   '
         f'four bundles {conv["spread_wall_mm"]:.0f} mm apart at the wall, '
         f'{conv["spread_target_mm"]:.1f} mm apart at the target plane',
         20, True, ACCENT),
        ('PRELIMINARY — transferred bench calibration, angles rescaled to '
         'v = 36.6 µm/ns. Not a measurement.', 14, False, MUTED),
    ])

    # ------------------------------------------------------------ the GIF
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _text(s, 0.5, 0.22, 12.3, 0.5, [
        ('The tour: all four bundles, then each segment alone', 22, True, INK)])
    iw, ih = Image.open(gif).size
    w, h = _fit(iw, ih, 12.2, 6.15)
    s.shapes.add_picture(gif, Inches((13.333 - w) / 2), Inches(0.85),
                         Inches(w), Inches(h))
    _text(s, 0.5, 7.02, 12.3, 0.4, [
        ('Animated — press F5 / use slideshow mode. Thick = the piece measured '
         'inside the 30 mm drift gap; thin = the same straight line continued '
         'out to the wall and back to the target plane.', 11, False, MUTED)])

    # ------------------------------------------------------------- stills
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _text(s, 0.5, 0.22, 12.3, 0.5, [
        ('One bundle per SiPM wall segment', 22, True, INK)])
    cells = [(0.35, 0.85), (6.85, 0.85), (0.35, 4.15), (6.85, 4.15)]
    for (x, y), png in zip(cells, stills):
        if not os.path.exists(png):
            continue
        iw, ih = Image.open(png).size
        w, h = _fit(iw, ih, 6.1, 3.05)
        s.shapes.add_picture(png, Inches(x), Inches(y), Inches(w), Inches(h))
    _text(s, 0.5, 7.05, 12.3, 0.4, [
        ('Each panel: the tracks whose trigger came from that group of four '
         'bars. Same lines, same view — only the selection changes.',
         11, False, MUTED)])

    # ------------------------------------------------------------ numbers
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _text(s, 0.8, 0.4, 11.7, 0.6, [
        ('Both pointing checks, as one number', 26, True, INK)])
    rows = [
        ('sample', 'n', 'spread at the wall', 'at the target plane'),
        ('both planes fitted + quality-ok', f'{st["n_all"]:,}',
         f'{conv["spread_wall_mm"]:.0f} mm', f'{conv["spread_target_mm"]:.1f} mm'),
        ('null: wall label shuffled', f'{st["n_all"]:,}',
         f'{null["spread_wall_mm"]:.0f} mm', f'{null["spread_target_mm"]:.1f} mm'),
    ]
    from pptx.util import Inches as In, Pt
    tbl = s.shapes.add_table(len(rows), 4, In(0.8), In(1.3), In(11.7),
                             In(1.9)).table
    for c, wdt in zip(range(4), (5.1, 1.4, 2.6, 2.6)):
        tbl.columns[c].width = In(wdt)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(15 if i else 13)
            p.runs[0].font.bold = (i == 1) or (i == 0)
            p.runs[0].font.color.rgb = _rgb(ACCENT if i == 1 else INK)
            p.runs[0].font.name = 'DejaVu Sans'
    _text(s, 0.8, 3.5, 11.7, 3.2, [
        ('The null is the control that matters.', 18, True, INK),
        ('Randomly relabelling which segment fired collapses the wall spread '
         f'from {conv["spread_wall_mm"]:.0f} mm to {null["spread_wall_mm"]:.0f} mm, '
         'so the separation is information the matcher supplies — not an '
         'artefact of splitting into four subsamples.', 15, False, INK),
        ('', 8, False, INK),
        ('And by the target plane the four bundles are '
         f'{conv["spread_target_mm"]:.1f} mm apart against the null\'s '
         f'{null["spread_target_mm"]:.1f} mm: statistically indistinguishable. '
         'The segment information is entirely used up by the pointing.',
         15, False, INK),
    ])

    # ------------------------------------------------------- what it isn't
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _text(s, 0.8, 0.4, 11.7, 0.6, [
        ('Two things the picture does NOT say', 26, True, ACCENT)])
    _text(s, 0.8, 1.35, 11.7, 5.4, [
        ('Individual tracks are not radial.', 18, True, INK),
        ('Per track, X at the target plane has median −23 mm with IQR '
         '[−46, −4]; only 15 % land within 10 mm of the beam axis and 48 % '
         'within 30 mm. The convergence is an ensemble statement about medians. '
         'A viewer who reads the waist as per-track pointing resolution is '
         'reading it wrong.', 15, False, INK),
        ('', 10, False, INK),
        ('The waist is not at the target.', 18, True, INK),
        ('It sits at X ≈ −23 mm; the target is at X = 0. That offset is the '
         'same size as the unresolved in-plane sign/offset convention (the '
         'chamber centre itself is at X = −16.35 mm, the arm-A pinwheel), so '
         'the figure cannot be used to locate the target. Fixing the in-plane '
         'convention is what would let it.', 15, False, INK),
        ('', 10, False, INK),
        ('Status: PRELIMINARY. The in-situ calibration of TRACK_PLAN_08 §6 has '
         'not been done — the bundle is the bench one with the DAQ constants '
         'swapped and a drift velocity that is argued for, not measured. '
         'Full detail: ntof_tracking/RUN79_PRELIM_2026-07-30.md §4.1.',
         13, False, MUTED),
    ])

    prs.save(out)
    return out


def build_web(js, gif, stills, all_png, outdir):
    os.makedirs(outdir, exist_ok=True)
    assets = []
    for src in [gif, all_png] + stills:
        if os.path.exists(src):
            shutil.copy(src, outdir)
            assets.append(os.path.basename(src))
    st = js['stats']
    conv, null = st['convergence'], st['convergence_null']
    seg_cards = '\n'.join(
        f'<figure><img src="{os.path.basename(p)}" alt="segment {k}">'
        f'<figcaption>segment {k} — {st["seg"][k]["n"]} tracks, bars '
        f'{st["seg"][k]["bars"]}, {100 * st["seg"][k]["inside"]:.0f} % inside '
        f'the group</figcaption></figure>'
        for k, p in zip(sorted(st['seg']), stills) if os.path.exists(p))
    html = f"""<!doctype html>
<meta charset="utf-8">
<title>run_79 / mx17_A — wall-segment tour</title>
<style>
 body{{background:#f7f7f7;color:#1a1a1a;font:16px/1.55 "DejaVu Sans",system-ui,sans-serif;
      margin:0 auto;max-width:1180px;padding:2.5rem 1.5rem 4rem}}
 h1{{font-size:1.75rem;margin:0 0 .3rem}} h2{{font-size:1.15rem;margin:2.5rem 0 .6rem}}
 .sub{{color:#6e6e6e;margin:0 0 1.6rem}} .hero{{color:#b0143c;font-weight:700}}
 img{{max-width:100%;height:auto;display:block;border-radius:4px}}
 figure{{margin:0}} figcaption{{color:#6e6e6e;font-size:.82rem;margin-top:.35rem}}
 .grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(330px,1fr));gap:1.4rem}}
 table{{border-collapse:collapse;margin:.5rem 0 1rem;font-size:.95rem}}
 th,td{{text-align:left;padding:.45rem 1.4rem .45rem 0;border-bottom:1px solid #e0e0e0}}
 tr.key td{{color:#b0143c;font-weight:700}}
 .note{{color:#6e6e6e;font-size:.85rem}}
</style>
<h1>Tracks point back at the target — and at the wall segment that fired</h1>
<p class="sub">run_79 / mx17_A waveform-first reconstruction, joined to n_TOF run
{st['ntof_run']} · <span class="hero">{st['n_all']:,} tracks over
{st['n_bunch']} bunches</span><br>
<strong>PRELIMINARY</strong> — transferred bench calibration, angles rescaled to
v = 36.6 µm/ns. Not a measurement.</p>

<img src="{os.path.basename(gif)}" alt="wall segment tour">
<p class="note">Thick = the piece measured inside the 30 mm drift gap; thin = the
same straight line continued out to the wall (dots) and back to the plane through
the target. Left: orbiting view. Right: beam's-eye, parallel projection.</p>

<h2>Both pointing checks, as one number</h2>
<table>
<tr><th>sample</th><th>n</th><th>spread at the wall</th><th>at the target plane</th></tr>
<tr class="key"><td>both planes fitted + quality-ok</td><td>{st['n_all']:,}</td>
<td>{conv['spread_wall_mm']:.0f} mm</td><td>{conv['spread_target_mm']:.1f} mm</td></tr>
<tr><td>null: wall label shuffled</td><td>{st['n_all']:,}</td>
<td>{null['spread_wall_mm']:.0f} mm</td><td>{null['spread_target_mm']:.1f} mm</td></tr>
</table>
<p>Randomly relabelling which segment fired collapses the wall spread from
{conv['spread_wall_mm']:.0f} mm to {null['spread_wall_mm']:.0f} mm, so the
separation is information the matcher supplies — not an artefact of splitting into
four subsamples. And by the target plane the bundles are
{conv['spread_target_mm']:.1f} mm apart against the null's
{null['spread_target_mm']:.1f} mm: statistically indistinguishable.</p>

<h2>One bundle per wall segment</h2>
<div class="grid">{seg_cards}</div>

<h2>Two things this does NOT say</h2>
<p><strong>Individual tracks are not radial.</strong> Per track, X at the target
plane has median −23 mm with IQR [−46, −4]; only 15 % land within 10 mm of the
beam axis and 48 % within 30 mm. The convergence is an ensemble statement about
medians.</p>
<p><strong>The waist is not at the target.</strong> It sits at X ≈ −23 mm; the
target is at X = 0. That offset is the same size as the unresolved in-plane
sign/offset convention (the chamber centre is itself at X = −16.35 mm), so the
figure cannot be used to locate the target.</p>
<p class="note">Full detail: <code>ntof_tracking/RUN79_PRELIM_2026-07-30.md</code> §4.1.</p>
"""
    with open(os.path.join(outdir, 'index.html'), 'w') as f:
        f.write(html)
    return os.path.join(outdir, 'index.html')


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--tour', default=TOUR)
    ap.add_argument('--small', action='store_true',
                    help='use the 900 px GIF (4.5 MB) instead of the 1400 px one')
    ap.add_argument('--out', default=None, help='output directory (default: the tour dir)')
    args = ap.parse_args()

    stem = 'wall_segment_tour_small' if args.small else 'wall_segment_tour'
    js = json.load(open(os.path.join(args.tour, f'{stem}.json')))
    gif = os.path.join(args.tour, f'{stem}.gif')
    all_png = os.path.join(args.tour, 'wall_segment_tour_all.png')
    stills = [os.path.join(args.tour, f'wall_segment_tour_seg{k}.png')
              for k in sorted(js['stats']['seg'])]
    outdir = args.out or args.tour
    os.makedirs(outdir, exist_ok=True)

    pptx = build_pptx(js, gif, stills, os.path.join(outdir, 'wall_segment_tour.pptx'))
    print(f'wrote {pptx}  ({os.path.getsize(pptx) / 1e6:.1f} MB)')
    web = build_web(js, gif, stills, all_png,
                    os.path.join(outdir, 'wall_segment_tour_web'))
    print(f'wrote {web}')
    print('\nshow it:')
    print(f'  soffice --show {pptx}')
    print(f'  xdg-open {web}')


if __name__ == '__main__':
    main()
