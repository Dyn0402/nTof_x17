#!/usr/bin/env python3
"""Convert slides/index.html into an editable PowerPoint deck.

Not a picture-per-slide export: every title, bullet, caption, table cell and
figure lands as a native PowerPoint shape you can click and retype.

The layout is not re-derived from the CSS -- it is *measured*.  The deck is
loaded in headless Chromium at exactly 1280x720, which is the size .deck
resolves to at 16:9 and therefore the size every cqw/cqh in the stylesheet is
computed against (the same reasoning make_pdf.sh relies on: @page is
13.333in x 7.5in = 1280 x 720 px, so paper, screen and slide share one set of
proportions).  Each slide is made .active in turn and every leaf element's
getBoundingClientRect() is read back.  PowerPoint then gets absolute boxes at
13.333in / 1280px = 9525 EMU per px.

What that buys: the browser does flexbox, grid, the container-query type scale
and line breaking, so the shapes land where Chromium put them rather than
where a hand-written CSS approximation guesses.

Fonts: the deck asks for Noto Sans Display / Noto Sans, and the weights it
uses (400/600/650/700/800) do not fit PowerPoint's two-faces-per-family limit.
`make_pptx_fonts.py` cuts the in-between weights out of the Noto variable
fonts as families of their own; install those first, or the measuring browser
and PowerPoint will disagree about how wide every line is.

    python make_pptx_fonts.py <dir with the Noto *-VF.ttf files>
    python to_pptx.py --assets <dir holding assets/img> --out talk.pptx

--base is for the case where the figures are not beside index.html (they are
gitignored renders, so on a fresh clone they usually are not); point it at a
file:// URL for a directory that has assets/img under it.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent

PX = 9525            # EMU per CSS px, at 13.333 in / 1280 px
DECK_W, DECK_H = 1280, 720

# .slide::before / ::after.  They are pseudo-elements, so they cannot be
# walked, and getComputedStyle returns `counter(slide)` unresolved -- so the
# text is rebuilt here using the same numbering make_pdf.sh uses.
FOOTER = "D. Neff · MPGD 2026, Prague · 3 September 2026"
TOTAL = "29"

# PowerPoint can only address two faces per family name -- regular and bold --
# but the deck uses font-weight 400, 600, 650, 700 and 800.  The weights in
# between therefore need family names of their own, the way Windows ships
# "Segoe UI Semibold" beside "Segoe UI"; scratch/mkfonts.py cuts them out of
# the Noto variable fonts.  This stylesheet pins the browser to exactly the
# same faces, so what is measured is what PowerPoint will lay out.
WEIGHT_CSS = """
<style id="pptx-weight-pins">
@font-face{font-family:"Noto Sans";font-weight:600;src:local("Noto Sans SemiBold")}
@font-face{font-family:"Noto Sans";font-weight:650;src:local("Noto Sans SemiBold")}
@font-face{font-family:"Noto Sans Display";font-weight:600;src:local("Noto Sans Display SemiBold")}
@font-face{font-family:"Noto Sans Display";font-weight:650;src:local("Noto Sans Display SemiBold")}
@font-face{font-family:"Noto Sans Display";font-weight:800;src:local("Noto Sans Display ExtraBold")}
</style>
"""


def face(run):
    """(family name, bold flag) for a measured run's weight and family."""
    wt = run.get("wt", 400)
    if run.get("mono"):
        return "Noto Sans Mono", wt >= 680
    stem = "Noto Sans Display" if run.get("disp") else "Noto Sans"
    if wt >= 750 and stem == "Noto Sans Display":
        return "Noto Sans Display ExtraBold", False
    if 550 <= wt < 680:
        return stem + " SemiBold", False
    return stem, wt >= 680

WALKER = r"""
(slideIndex) => {
  const deck = document.querySelector('.deck');
  const D = deck.getBoundingClientRect();
  const sl = document.querySelectorAll('section.slide')[slideIndex];

  // Tags that live *inside* a run of text rather than forming a box of their
  // own.  <tag> and <n> are the deck's own inline idioms.
  const INLINE = new Set(['B','I','EM','STRONG','CODE','SPAN','SUB','SUP','BR',
                          'A','SMALL','U','MARK','TAG','N','WBR','ABBR','CITE']);

  const rect = el => {
    const r = el.getBoundingClientRect();
    return {x: r.left - D.left, y: r.top - D.top, w: r.width, h: r.height};
  };
  const vis = el => {
    const cs = getComputedStyle(el);
    if (cs.display === 'none' || cs.visibility === 'hidden') return false;
    if (parseFloat(cs.opacity) < 0.02) return false;
    return true;
  };
  const paint = c => c && c !== 'rgba(0, 0, 0, 0)' && c !== 'transparent';

  // Flatten an element's inline descendants into styled runs.
  const runsOf = el => {
    const out = [];
    const base = getComputedStyle(el);
    (function walk(n, st) {
      if (n.nodeType === 3) {
        if (n.nodeValue) out.push(Object.assign({t: n.nodeValue}, st));
        return;
      }
      if (n.nodeType !== 1) return;
      if (n.tagName === 'BR') { out.push({t: '\n', br: true}); return; }
      const cs = getComputedStyle(n);
      const st2 = {
        wt: parseInt(cs.fontWeight, 10) || 400,
        i: st.i || cs.fontStyle === 'italic',
        disp: /Display/.test(cs.fontFamily),
        mono: st.mono || /mono/i.test(cs.fontFamily),
        sub: st.sub || n.tagName === 'SUB',
        sup: st.sup || n.tagName === 'SUP',
        color: cs.color,
        fs: parseFloat(cs.fontSize),
      };
      for (const c of n.childNodes) walk(c, st2);
    })(el, {wt: parseInt(base.fontWeight, 10) || 400,
            i: base.fontStyle === 'italic',
            disp: /Display/.test(base.fontFamily),
            mono: /mono/i.test(base.fontFamily),
            sub: false, sup: false,
            color: base.color, fs: parseFloat(base.fontSize)});
    return out;
  };

  // Whether an element's text can collapse into ONE PowerPoint text box.
  // Being built only from inline tags is not enough: a flex or grid container
  // positions even <span>s as columns with a gap between them (.outline-item
  // does exactly that), and flattening would run the number into the title.
  // Anything whose children carry a box of their own is walked instead.
  const textish = el => {
    if (!el.textContent.trim()) return false;
    const cs = getComputedStyle(el);
    if (/flex|grid/.test(cs.display)) return false;
    for (const c of el.children) {
      if (!INLINE.has(c.tagName)) return false;
      if (getComputedStyle(c).display !== 'inline') return false;
    }
    return true;
  };

  // Text sitting directly inside a container that had to be walked rather
  // than flattened.  A Range gives it the box it actually occupies.
  const looseText = (el, node) => {
    const rng = document.createRange();
    rng.selectNode(node);
    const r = rng.getBoundingClientRect();
    rng.detach();
    if (!(r.width > 0 && r.height > 0)) return null;
    const cs = getComputedStyle(el);
    return {
      kind: 'text',
      rect: {x: r.left - D.left, y: r.top - D.top, w: r.width, h: r.height},
      runs: [{t: node.nodeValue, wt: parseInt(cs.fontWeight, 10) || 400,
              i: cs.fontStyle === 'italic', disp: /Display/.test(cs.fontFamily),
              mono: /mono/i.test(cs.fontFamily), color: cs.color,
              fs: parseFloat(cs.fontSize)}],
      align: cs.textAlign, fs: parseFloat(cs.fontSize), color: cs.color,
      lh: cs.lineHeight === 'normal' ? null : parseFloat(cs.lineHeight),
      ls: parseFloat(cs.letterSpacing) || 0,
      transform: cs.textTransform, cls: '',
    };
  };

  // A box worth drawing under the text: a fill, or a visible border.  The
  // deck uses these for the bar chart, callouts and hairlines.
  const boxOf = el => {
    const cs = getComputedStyle(el);
    const sides = ['Top', 'Right', 'Bottom', 'Left'];
    const bw = sides.map(s => parseFloat(cs['border' + s + 'Width']) || 0);
    if (!paint(cs.backgroundColor) && !bw.some(w => w > 0)) return null;
    return {
      fill: paint(cs.backgroundColor) ? cs.backgroundColor : null,
      radius: parseFloat(cs.borderTopLeftRadius) || 0,
      borders: sides.map((s, k) => bw[k] > 0
        ? {side: s.toLowerCase(), w: bw[k], color: cs['border' + s + 'Color']}
        : null).filter(Boolean),
    };
  };

  const tableOf = el => {
    const rows = [];
    for (const tr of el.querySelectorAll('tr')) {
      const cells = [];
      for (const td of tr.children) {
        const cs = getComputedStyle(td);
        cells.push({
          head: td.tagName === 'TH', rect: rect(td), runs: runsOf(td),
          align: cs.textAlign, fs: parseFloat(cs.fontSize),
          colspan: td.colSpan || 1,
          lh: cs.lineHeight === 'normal' ? null : parseFloat(cs.lineHeight),
          pad: ['Top', 'Right', 'Bottom', 'Left'].map(
                 s => parseFloat(cs['padding' + s]) || 0),
          borderTop: parseFloat(cs.borderTopWidth) || 0,
          borderBottom: parseFloat(cs.borderBottomWidth) || 0,
        });
      }
      rows.push({rect: rect(tr), cells});
    }
    return {rows};
  };

  const items = [];

  (function walk(el) {
    if (!vis(el)) return;
    const tag = el.tagName;

    if (tag === 'IMG') {
      const r = rect(el);
      const cs = getComputedStyle(el);
      // object-fit:contain letterboxes inside the element box; recover the
      // rectangle the pixels actually occupy so nothing is stretched.
      let box = r;
      if (cs.objectFit === 'contain' && el.naturalWidth) {
        const s = Math.min(r.w / el.naturalWidth, r.h / el.naturalHeight);
        const w = el.naturalWidth * s, h = el.naturalHeight * s;
        box = {x: r.x + (r.w - w) / 2, y: r.y + (r.h - h) / 2, w: w, h: h};
      }
      items.push({kind: 'image', rect: box, src: el.getAttribute('src'),
                  alt: el.getAttribute('alt') || ''});
      return;
    }
    if (tag === 'TABLE') {
      items.push(Object.assign({kind: 'table', rect: rect(el)}, tableOf(el)));
      return;
    }
    if (tag.toLowerCase() === 'svg') {
      items.push({kind: 'svg', rect: rect(el),
                  id: el.getAttribute('data-pptx-id')});
      return;
    }

    const box = boxOf(el);
    if (box && el !== sl) {
      items.push(Object.assign({kind: 'rect', rect: rect(el)}, box));
    }

    // Decorative pseudo-elements: the kicker rule, the title rule, the two
    // divider rules.  Their used width/height/background resolve fine even
    // though a counter() content string does not -- but they have no box of
    // their own to measure, so it is reconstructed from the parent's, and the
    // parent's text is inset by however much the rule displaces it.
    let insetL = 0, insetT = 0;
    for (const pe of ['::before', '::after']) {
      const cs = getComputedStyle(el, pe);
      if (!cs.content || cs.content === 'none' || cs.content === 'normal') continue;
      if (!paint(cs.backgroundColor)) continue;
      const w = parseFloat(cs.width), h = parseFloat(cs.height);
      if (!(w > 0 && h > 0)) continue;
      const p = rect(el);
      const ml = parseFloat(cs.marginLeft) || 0, mr = parseFloat(cs.marginRight) || 0;
      const mt = parseFloat(cs.marginTop) || 0, mb = parseFloat(cs.marginBottom) || 0;
      const pcs = getComputedStyle(el);
      const inline = /inline/.test(cs.display);
      let box;
      if (inline) {
        // Sits on the first line, vertically centred in it.
        const lh = pcs.lineHeight === 'normal'
          ? parseFloat(pcs.fontSize) * 1.2 : parseFloat(pcs.lineHeight);
        box = {x: p.x + insetL + ml, y: p.y + (lh - h) / 2, w: w, h: h};
        if (pe === '::before') insetL += w + ml + mr;
      } else {
        box = {x: p.x + ml, y: p.y + insetT + mt, w: w, h: h};
        if (pe === '::before') insetT += h + mt + mb;
      }
      items.push({kind: 'pseudo', pe: pe, rect: box, fill: cs.backgroundColor});
    }
    el.__pptxInset = {l: insetL, t: insetT};

    if (tag === 'UL' || tag === 'OL') {
      const lis = [];
      for (const li of el.children) {
        if (!vis(li)) continue;
        const mk = getComputedStyle(li, '::marker');
        const lcs = getComputedStyle(li);
        lis.push({rect: rect(li), runs: runsOf(li),
                  marker: (mk.content || '').replace(/^"|"$/g, ''),
                  markerColor: mk.color,
                  fs: parseFloat(lcs.fontSize),
                  align: lcs.textAlign,
                  lh: lcs.lineHeight === 'normal' ? null
                      : parseFloat(lcs.lineHeight)});
      }
      const cs = getComputedStyle(el);
      items.push({kind: 'list', rect: rect(el), items: lis,
                  lh: cs.lineHeight === 'normal' ? null : parseFloat(cs.lineHeight),
                  align: cs.textAlign});
      return;
    }

    if (textish(el)) {
      const cs = getComputedStyle(el);
      const ins = el.__pptxInset || {l: 0, t: 0};
      const r = rect(el);
      // getBoundingClientRect is the border box.  A padded block (.callout is
      // the visible case) must hand PowerPoint its CONTENT box, or the text
      // starts flush against the panel edge.
      const pl = parseFloat(cs.paddingLeft) || 0;
      const pr = parseFloat(cs.paddingRight) || 0;
      const ptp = parseFloat(cs.paddingTop) || 0;
      const pb = parseFloat(cs.paddingBottom) || 0;
      items.push({
        kind: 'text',
        rect: {x: r.x + pl + ins.l, y: r.y + ptp + ins.t,
               w: r.w - pl - pr - ins.l, h: r.h - ptp - pb - ins.t},
        runs: runsOf(el),
        align: cs.textAlign,
        lh: cs.lineHeight === 'normal' ? null : parseFloat(cs.lineHeight),
        fs: parseFloat(cs.fontSize), color: cs.color,
        ls: parseFloat(cs.letterSpacing) || 0,
        transform: cs.textTransform,
        cls: typeof el.className === 'string' ? el.className : '',
      });
      return;
    }

    for (const c of el.childNodes) {
      if (c.nodeType === 1) walk(c);
      else if (c.nodeType === 3 && c.nodeValue.trim()) {
        const t = looseText(el, c);
        if (t) items.push(t);
      }
    }
  })(sl);

  const g = cs => (!cs.content || cs.content === 'none') ? null : {
    fs: parseFloat(cs.fontSize), color: cs.color,
    left: parseFloat(cs.left), right: parseFloat(cs.right),
    bottom: parseFloat(cs.bottom), ls: parseFloat(cs.letterSpacing) || 0,
    transform: cs.textTransform, wt: parseInt(cs.fontWeight, 10) || 400,
  };
  return {
    classes: sl.className, frame: sl.getAttribute('data-frame'),
    items: items,
    footer: {before: g(getComputedStyle(sl, '::before')),
             after: g(getComputedStyle(sl, '::after'))},
  };
}
"""


def measure(index_html, base_href, out_json, channel):
    from playwright.sync_api import sync_playwright

    html = index_html.read_text(encoding="utf-8")

    # Give every inline <svg> a handle so it can be screenshotted by id.
    n = [0]

    def tag_svg(m):
        n[0] += 1
        return m.group(0)[:-1] + ' data-pptx-id="svg%d">' % n[0]

    html = re.sub(r"<svg\b[^>]*>", tag_svg, html)
    html = WEIGHT_CSS + html
    if base_href:
        html = '<base href="%s">\n' % base_href + html

    work = out_json.parent / "_measure.html"
    work.write_text(html, encoding="utf-8")

    shots = out_json.parent / "svg"
    shots.mkdir(parents=True, exist_ok=True)

    slides = []
    with sync_playwright() as p:
        kw = {"channel": channel} if channel else {}
        br = p.chromium.launch(**kw)
        pg = br.new_page(viewport={"width": DECK_W, "height": DECK_H},
                         device_scale_factor=2)
        pg.goto(work.as_uri())
        pg.wait_for_load_state("networkidle")

        total = pg.eval_on_selector_all("section.slide", "e => e.length")
        print("measuring %d slides" % total, file=sys.stderr)

        for i in range(total):
            pg.evaluate(
                """(i) => {
                    const all = document.querySelectorAll('section.slide');
                    all.forEach(s => s.classList.remove('active'));
                    all[i].classList.add('active');
                }""", i)
            rec = pg.evaluate(WALKER, i)
            rec["index"] = i
            # Inline SVGs are re-drawn as pictures: there are four in the deck
            # and none is text the speaker will want to edit.
            for it in rec["items"]:
                if it["kind"] == "svg" and it.get("id"):
                    f = shots / ("%s.png" % it["id"])
                    pg.locator("[data-pptx-id='%s']" % it["id"]).screenshot(
                        path=str(f))
                    it["file"] = str(f)
            slides.append(rec)
            if (i + 1) % 20 == 0:
                print("  %d/%d" % (i + 1, total), file=sys.stderr)
        br.close()

    out_json.write_text(json.dumps(slides), encoding="utf-8")
    print("wrote %s (%d slides)" % (out_json, len(slides)), file=sys.stderr)
    return slides


# --------------------------------------------------------------------------
# Emitting the .pptx
# --------------------------------------------------------------------------

RGB_RE = re.compile(r"rgba?\(([^)]+)\)")

# hhea ascent + descent, as a fraction of the em.  Every Noto face installed
# for this deck reports the same 1.069 + 0.293, so one constant covers them
# all; re-measure if the deck ever moves to another family.
EM_HEIGHT = 1.362


def rgb(css):
    """'rgb(35, 55, 59)' -> RGBColor, or None if fully transparent."""
    from pptx.dml.color import RGBColor
    if not css:
        return None
    m = RGB_RE.match(css.strip())
    if not m:
        return None
    parts = [p.strip() for p in m.group(1).replace("/", " ").split(",")]
    vals = [float(p) for p in parts[:3]]
    if len(parts) > 3 and float(parts[3]) < 0.02:
        return None
    return RGBColor(int(vals[0]), int(vals[1]), int(vals[2]))


def emu(px):
    return int(round(px * PX))


def pt(px):
    """CSS px -> points.  1 CSS px is 1/96 in here, so exactly 0.75 pt."""
    from pptx.util import Pt
    return Pt(px * 0.75)


def normalise(runs):
    """Collapse HTML whitespace across a run list, keeping the run styling.

    Done across the whole list rather than run by run, so that the space
    between '<b>foo</b> bar' survives exactly once and the indentation
    newlines in the source do not become spaces at the start of a line.
    """
    out = []
    prev_space = True  # so leading whitespace is dropped
    for r in runs:
        if r.get("br"):
            out.append(dict(r, t="\n"))
            prev_space = True
            continue
        t = re.sub(r"\s+", " ", r.get("t", ""))
        if not t:
            continue
        if prev_space and t.startswith(" "):
            t = t[1:]
        if not t:
            continue
        prev_space = t.endswith(" ")
        out.append(dict(r, t=t))
    while out and not out[-1].get("br") and out[-1]["t"].strip() == "":
        out.pop()
    if out and not out[-1].get("br"):
        out[-1] = dict(out[-1], t=out[-1]["t"].rstrip())
    return out


def style_run(run, spec, caps=False, ls=0.0):
    from pptx.util import Pt
    fam, bold = face(spec)
    f = run.font
    f.name = fam
    f.bold = bold
    f.italic = bool(spec.get("i"))
    f.size = pt(spec.get("fs", 16))
    c = rgb(spec.get("color"))
    if c is not None:
        f.color.rgb = c
    rPr = f._rPr
    if caps:
        rPr.set("cap", "all")
    if ls:
        rPr.set("spc", str(int(round(ls * 0.75 * 100))))
    if spec.get("sup"):
        rPr.set("baseline", "30000")
    elif spec.get("sub"):
        rPr.set("baseline", "-25000")


ALIGN = {}


def _align(name):
    from pptx.enum.text import PP_ALIGN
    if not ALIGN:
        ALIGN.update({"left": PP_ALIGN.LEFT, "start": PP_ALIGN.LEFT,
                      "right": PP_ALIGN.RIGHT, "end": PP_ALIGN.RIGHT,
                      "center": PP_ALIGN.CENTER, "justify": PP_ALIGN.JUSTIFY})
    return ALIGN.get(name)


def add_text(slide, item, runs=None, bullet=None):
    """One HTML text element -> one PowerPoint text box, positioned as measured."""
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.util import Pt

    runs = normalise(runs if runs is not None else item["runs"])
    if not runs:
        return None
    r = item["rect"]
    # Half-leading.  CSS centres a line's glyphs in its line box, so a
    # line-height tighter than the font's own ascent+descent crops evenly top
    # and bottom.  PowerPoint instead hangs the text from the top of the line,
    # which lifts every line by that same half-leading -- about 3 px at title
    # sizes here.  Push the box back down by it.
    lh = item.get("lh")
    fs = item.get("fs") or max((s.get("fs", 16) for s in runs), default=16)
    lift = max(0.0, (EM_HEIGHT * fs - lh) / 2) if lh else 0.0

    # A little slack on the right and bottom: PowerPoint's line breaking is
    # not bit-identical to Chromium's, and a box sized to the last pixel wraps
    # a word that fitted in the browser.  The box is transparent, so growing
    # it cannot show.
    box = slide.shapes.add_textbox(emu(r["x"]), emu(r["y"] + lift),
                                   emu(r["w"] + 2), emu(r["h"] + 4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    caps = item.get("transform") == "uppercase"
    ls = item.get("ls", 0) or 0
    lh = item.get("lh")

    paras = [[]]
    for run in runs:
        if run.get("br"):
            paras.append([])
        else:
            paras[-1].append(run)

    for k, prun in enumerate(paras):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        al = _align(item.get("align"))
        if al is not None:
            p.alignment = al
        if lh:
            p.line_spacing = pt(lh)
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if bullet:
            _bullet(p, bullet)
        for spec in prun:
            run = p.add_run()
            run.text = spec["t"]
            style_run(run, spec, caps=caps, ls=ls)
        if not prun:
            p.add_run().text = ""
    return box


def _bullet(p, bullet):
    """Give a paragraph a real PowerPoint bullet, not a typed-in character."""
    from pptx.oxml.ns import qn
    pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
    ch, colour, indent = bullet
    pPr.set("indent", str(-indent))
    pPr.set("marL", str(indent))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    if colour is not None:
        clr = pPr.makeelement(qn("a:buClr"), {})
        srgb = pPr.makeelement(qn("a:srgbClr"), {"val": str(colour)})
        clr.append(srgb)
        pPr.append(clr)
    bu = pPr.makeelement(qn("a:buChar"), {"char": ch})
    pPr.append(bu)


def add_rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    if w <= 0 or h <= 0:
        return None
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y),
                                max(emu(w), 1), max(emu(h), 1))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.text_frame.text = ""
    return sh


def add_table(slide, item):
    from pptx.oxml.ns import qn
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Pt

    rows = item["rows"]
    if not rows:
        return
    ncols = max(sum(c["colspan"] for c in r["cells"]) for r in rows)
    r0 = item["rect"]
    shape = slide.shapes.add_table(len(rows), ncols, emu(r0["x"]), emu(r0["y"]),
                                   emu(r0["w"]), emu(r0["h"]))
    tbl = shape.table

    # Strip the default banded blue look; the deck's tables are booktabs
    # rules on white, which are drawn separately as thin rectangles below.
    tbl.first_row = False
    tbl.horz_banding = False
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    for e in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(e)
    sid = tblPr.makeelement(qn("a:tableStyleId"), {})
    sid.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid
    tblPr.append(sid)

    widths = [None] * ncols
    for row in rows:
        col = 0
        for c in row["cells"]:
            if c["colspan"] == 1 and widths[col] is None:
                widths[col] = c["rect"]["w"]
            col += c["colspan"]
    span = sum(w for w in widths if w)
    for k, w in enumerate(widths):
        tbl.columns[k].width = emu(w if w else
                                   max(r0["w"] - span, 0) / max(
                                       sum(1 for x in widths if not x), 1))
    for k, row in enumerate(rows):
        tbl.rows[k].height = emu(row["rect"]["h"])

    for k, row in enumerate(rows):
        col = 0
        for c in row["cells"]:
            cell = tbl.cell(k, col)
            if c["colspan"] > 1:
                cell.merge(tbl.cell(k, col + c["colspan"] - 1))
            # HTML pads a cell from the top; a PowerPoint cell would centre
            # its text in whatever height the row ends up with, which drifts
            # further with every row.  Anchor to the top and reproduce the
            # measured padding as the cell's own margins instead.
            pad = c.get("pad") or [0, 0, 0, 0]
            lift = 0.0
            if c.get("lh"):
                lift = max(0.0, (EM_HEIGHT * c.get("fs", 16) - c["lh"]) / 2)
            cell.margin_top = emu(pad[0] + lift)
            cell.margin_right = emu(pad[1])
            cell.margin_bottom = emu(pad[2])
            cell.margin_left = emu(pad[3])
            cell.vertical_anchor = MSO_ANCHOR.TOP
            cell.fill.background()
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            al = _align(c.get("align"))
            if al is not None:
                p.alignment = al
            if c.get("lh"):
                p.line_spacing = pt(c["lh"])
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            for spec in normalise(c["runs"]):
                if spec.get("br"):
                    p = tf.add_paragraph()
                    if al is not None:
                        p.alignment = al
                    if c.get("lh"):
                        p.line_spacing = pt(c["lh"])
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    continue
                run = p.add_run()
                run.text = spec["t"]
                style_run(run, spec)
            col += c["colspan"]

    # The horizontal rules the design actually uses.
    for row in rows:
        for c in row["cells"]:
            rr = c["rect"]
            for w, yy in ((c.get("borderTop", 0), rr["y"]),
                          (c.get("borderBottom", 0), rr["y"] + rr["h"] - c.get("borderBottom", 0))):
                if w:
                    add_rect(slide, rr["x"], yy, rr["w"], w, rgb("rgb(35,55,59)"))


def build(slides, out_pptx, asset_root, notes=None):
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(DECK_W * PX)
    prs.slide_height = Emu(DECK_H * PX)
    blank = prs.slide_layouts[6]

    paper = rgb("rgb(255,255,255)")
    s_no = 0
    missing = []

    for rec in slides:
        classes = rec["classes"].split()
        if "bcont" not in classes:
            s_no += 1
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, DECK_W, DECK_H, paper)

        for item in rec["items"]:
            k = item["kind"]
            if k == "rect":
                r = item["rect"]
                if item.get("fill"):
                    add_rect(slide, r["x"], r["y"], r["w"], r["h"],
                             rgb(item["fill"]))
                for b in item.get("borders", []):
                    c = rgb(b["color"])
                    w = b["w"]
                    if b["side"] == "top":
                        add_rect(slide, r["x"], r["y"], r["w"], w, c)
                    elif b["side"] == "bottom":
                        add_rect(slide, r["x"], r["y"] + r["h"] - w, r["w"], w, c)
                    elif b["side"] == "left":
                        add_rect(slide, r["x"], r["y"], w, r["h"], c)
                    else:
                        add_rect(slide, r["x"] + r["w"] - w, r["y"], w, r["h"], c)
            elif k == "pseudo":
                r = item["rect"]
                add_rect(slide, r["x"], r["y"], r["w"], r["h"], rgb(item["fill"]))
            elif k == "image":
                r = item["rect"]
                src = item["src"]
                path = (asset_root / src) if asset_root else Path(src)
                if not path.exists():
                    missing.append(src)
                    continue
                pic = slide.shapes.add_picture(str(path), emu(r["x"]), emu(r["y"]),
                                               emu(r["w"]), emu(r["h"]))
                if item.get("alt"):
                    pic._element._nvXxPr.cNvPr.set("descr", item["alt"][:1000])
            elif k == "svg":
                r = item["rect"]
                if item.get("file") and Path(item["file"]).exists():
                    slide.shapes.add_picture(item["file"], emu(r["x"]), emu(r["y"]),
                                             emu(r["w"]), emu(r["h"]))
            elif k == "table":
                add_table(slide, item)
            elif k == "list":
                for li in item["items"]:
                    mark = (li.get("marker") or "▸").strip() or "▸"
                    col = rgb(li.get("markerColor"))
                    # The measured rect starts where the TEXT starts -- an
                    # outside marker hangs in the ul's padding.  PowerPoint
                    # instead indents the text away from the box's own left
                    # edge, so the box has to start back at the marker or
                    # every bullet loses a marker's width of column and wraps
                    # a word early.
                    ind = round(li.get("fs", 16) * 1.15)
                    r = li["rect"]
                    li = dict(li, rect={"x": r["x"] - ind, "y": r["y"],
                                        "w": r["w"] + ind, "h": r["h"]})
                    add_text(slide, li, runs=li["runs"],
                             bullet=(mark[0], col, int(ind) * PX))
            elif k == "text":
                add_text(slide, item)

        # The running footer and the slide number, which are pseudo-elements
        # in CSS and so have to be rebuilt rather than measured.
        fb, fa = rec["footer"]["before"], rec["footer"]["after"]
        if fb:
            _footer(slide, fb, FOOTER, DECK_W - fb["left"] - fb["right"],
                    fb["left"], "left")
        if fa:
            num = "%d" % s_no
            if rec.get("frame"):
                num += "." + rec["frame"]
            _footer(slide, fa, "%s / %s" % (num, TOTAL),
                    DECK_W - fa["left"] - fa["right"], fa["left"], "right")

        if notes and rec["index"] in notes:
            slide.notes_slide.notes_text_frame.text = notes[rec["index"]]

    prs.save(str(out_pptx))
    return missing


def _footer(slide, cs, text, width, left, align):
    h = cs["fs"] * 1.6
    item = {"rect": {"x": left, "y": DECK_H - cs["bottom"] - h, "w": width, "h": h},
            "align": align, "ls": cs.get("ls", 0), "lh": None, "fs": cs["fs"],
            "transform": cs.get("transform", "none")}
    add_text(slide, item, runs=[{"t": text, "wt": cs.get("wt", 400),
                                 "fs": cs["fs"], "color": cs["color"]}])


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--index", default=str(HERE / "index.html"))
    ap.add_argument("--base", default=None,
                    help="<base href> for assets not beside index.html")
    ap.add_argument("--json", default=None)
    ap.add_argument("--channel", default="msedge",
                    help="playwright channel ('' for bundled chromium)")
    ap.add_argument("--assets", default=None,
                    help="directory the img src paths are relative to")
    ap.add_argument("--out", default=None, help="output .pptx")
    ap.add_argument("--reuse", action="store_true",
                    help="skip the browser pass and use the existing --json")
    a = ap.parse_args()

    out = Path(a.json) if a.json else HERE / "_pptx_measure.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    if a.reuse and out.exists():
        slides = json.loads(out.read_text(encoding="utf-8"))
        print("reusing %s (%d slides)" % (out, len(slides)), file=sys.stderr)
    else:
        slides = measure(Path(a.index), a.base, out, a.channel)

    pptx = Path(a.out) if a.out else HERE / "mpgd26_talk.pptx"
    root = Path(a.assets) if a.assets else Path(a.index).resolve().parent
    missing = build(slides, pptx, root)
    if missing:
        print("WARNING: %d image(s) not found, e.g. %s"
              % (len(missing), missing[0]), file=sys.stderr)
    print("wrote %s (%d slides)" % (pptx, len(slides)), file=sys.stderr)


if __name__ == "__main__":
    main()
