"""
build_pptx.py — Rebuild the MIP slides as a .pptx that mirrors Dylan's usual
deck formatting, by using the downloaded deck itself as the template (master,
layouts, theme fonts/colors inherited; its slides removed).

Usage: python build_pptx.py
Output: mip_slides.pptx (the beamer mip_slides.pdf remains untouched)
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

TEMPLATE = Path.home() / 'Downloads' / 'X17 Overview 7_13.pptx'
FIG = Path(__file__).parent.parent / 'figures'
OUT = Path(__file__).parent / 'mip_slides.pptx'

L_TITLE, L_BODY, L_TWOCOL, L_ONLY = 0, 2, 3, 5

prs = Presentation(TEMPLATE)
# drop template slides AND their package parts, keep master/layouts
from pptx.oxml.ns import qn
xml_slides = prs.slides._sldIdLst
for sld in list(xml_slides):
    prs.part.drop_rel(sld.get(qn('r:id')))
    xml_slides.remove(sld)

layouts = prs.slide_masters[0].slide_layouts
SW, SH = Inches(13.33), Inches(7.5)


def add(layout_i):
    return prs.slides.add_slide(layouts[layout_i])


def set_bullets(ph, items):
    tf = ph.text_frame
    tf.clear()
    for i, (lvl, txt) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.level = lvl


def add_picture(slide, path, top_in=1.15, max_w_in=12.5, max_h_in=5.3,
                center_x_in=13.33 / 2, left_in=None):
    w0, h0 = Image.open(path).size
    scale = min(Inches(max_w_in) / w0, Inches(max_h_in) / h0)
    w, h = int(w0 * scale), int(h0 * scale)
    left = Inches(left_in) if left_in is not None else int(Inches(center_x_in) - w / 2)
    return slide.shapes.add_picture(str(path), left, Inches(top_in), w, h)


def add_caption(slide, text, top_in=6.6, size=11):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top_in), Inches(12.1), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    return tb


def two_col(title, fig_path, bullets, fig_frac=0.58):
    s = add(L_TWOCOL)
    s.shapes.title.text = title
    left_ph = s.placeholders[1]
    right_ph = s.placeholders[2]
    # image where the left body placeholder sits (slightly enlarged)
    lft, top = left_ph.left, left_ph.top
    w0, h0 = Image.open(fig_path).size
    max_w, max_h = int(SW * fig_frac), int(Inches(5.4))
    scale = min(max_w / w0, max_h / h0)
    s.shapes.add_picture(str(fig_path), Inches(0.25), top,
                         int(w0 * scale), int(h0 * scale))
    left_ph._element.getparent().remove(left_ph._element)
    set_bullets(right_ph, bullets)
    return s


# ── 1. title ────────────────────────────────────────────────────────────────
s = add(L_TITLE)
s.shapes.title.text = 'SiPM Wall MIP Calibration from Wall–Plastic Coincidences'
s.placeholders[1].text = ('n_TOF EAR2 X17 — run 224404 (first beam data, July 15)\n'
                          'Dylan Neff — July 17, 2026')

# ── 2. data & selection ─────────────────────────────────────────────────────
s = add(L_BODY)
s.shapes.title.text = 'Data & selection'
set_bullets(s.placeholders[1], [
    (0, 'Official processed file, run 224404: 3792 bunches (1471 dedicated), ~300M hits'),
    (0, 'Per-channel ADC→mV from DAQ settings (≈0.0305 mV/count)'),
    (0, 'SiPM walls delivered NO data for 46% of beam-on bunches (643–2212) — '
        'plastics unaffected; hardware cause TBD with DAQ crew'),
    (0, 'All results: post-recovery bunches, hits >0.1 ms after γ-flash '
        '(highest coincidence purity)'),
])

# ── 3. method ───────────────────────────────────────────────────────────────
s = add(L_ONLY)
s.shapes.title.text = 'Method: coincidences with sideband subtraction'
add_picture(s, FIG / '05_sideband_diagram/sideband_method.png', max_h_in=5.6)

# ── 4. dt matrix ────────────────────────────────────────────────────────────
s = add(L_ONLY)
s.shapes.title.text = 'Every arm pairs wall ↔ plastic cleanly'
add_picture(s, FIG / '02_coincidence/dt_matrix.png', max_h_in=5.2)
add_caption(s, 'Peaks at −8 to −18 ns, rms 1.5–3.5 ns; per-channel '
               'offsets now a stored calibration (calib/time_offsets_run224404.json).')

# ── 5. geometry ─────────────────────────────────────────────────────────────
s = add(L_ONLY)
s.shapes.title.text = 'Geometry confirmed in data'
add_picture(s, FIG / '06_07_geometry_mip/wall_geometry_matrices.png', max_h_in=5.0)
add_caption(s, 'Top/bottom blocks (1,2)(3,4)(5,6)(7,8) per 4-bar group; groups ordered '
               'left→right across the two bars. Corner hotspots ⇒ A–D and '
               'B–C are ADJACENT arms. Odd=top is an assumption — need cabling sheet.')

# ── 6. headline MIP ─────────────────────────────────────────────────────────
s = add(L_ONLY)
s.shapes.title.text = 'Headline: SiPM-wall MIP peaks on arms B and C'
add_picture(s, FIG / '06_07_geometry_mip/mip_wall_spectra_linear.png', max_h_in=4.9)
add_caption(s, 'Sideband-subtracted coincidence spectra. B/C: clean MIP peaks, '
               '29–39 mV, all 16 channels ⇒ per-channel calibration. '
               'A/D: no MIP population (later slides).', top_in=6.3, size=13)

# ── 7. inclusive vs coincidence ─────────────────────────────────────────────
s = add(L_ONLY)
s.shapes.title.text = 'Coincidence selection performs exactly as intended (B/C)'
add_picture(s, FIG / '08_inclusive_vs_coinc/inclusive_vs_coinc_linear.png', max_h_in=4.9)
add_caption(s, 'Physics note: the plastic is at the BACK of the stack ⇒ the '
               'coincidence is a through-going (MIP) selection for the wall, but only '
               'a hit selection for the plastic — plastic spectra are NOT MIP spectra.',
            top_in=6.3, size=13)

# ── 8. A/D problem ──────────────────────────────────────────────────────────
two_col('The A/D problem — cause still open', FIG / '06_07_geometry_mip/mip_wall_spectra_linear.png', [
    (0, 'A/D: 6–8× fewer true coincidences; coincident spectrum just '
        'traces the inclusive shape — no MIP bump'),
    (0, 'Hypotheses: plastic response misses MIPs / wall SiPM gain low / '
        'geometry, absorption or flux composition'),
    (0, 'Inclusive plastic medians agree within ±30% across all 8 PMTs '
        '(DL highest) ⇒ large plastic-gain deficit unlikely'),
    (0, 'Wall-only top–bottom check (backup slide): SiPM response healthy on '
        'all four arms — MIP-scale deposits present everywhere'),
    (1, 'but top & bottom view the SAME bars, so a single gamma/Compton '
        'deposit also fires both ends — this does NOT prove a through-going '
        'flux, nor pin the blame on the plastic'),
    (0, 'Discriminating follow-ups: pre-flash cosmic sample (guaranteed '
        'through-going muons), LIQ readout, plastic HV scan'),
])

# ── 10. PMT gain table ──────────────────────────────────────────────────────
s = add(L_ONLY)
s.shapes.title.text = 'Plastic PMT relative gains & HV suggestions'
rows = [('PMT', 'HV [V]', 'med_inc [mV]', 'rel. gain', 'V (equalize)', 'V (2×)'),
        ('PSSA1 (AL)', '1325', '10.4', '0.74', '1384', '1528'),
        ('PSSA2 (AR)', '1275', '13.2', '0.93', '1288', '1422'),
        ('PSSB1 (BL)', '1325', '13.1', '0.93', '1339', '1478'),
        ('PSSB2 (BR)', '1300', '14.9', '1.05', '1290', '1425'),
        ('PSSC1 (CL)', '1300', '14.5', '1.03', '1295', '1429'),
        ('PSSC2 (CR)', '1300', '15.3', '1.08', '1285', '1419'),
        ('PSSD1 (DL)', '1300', '17.8', '1.26', '1258', '1389'),
        ('PSSD2 (DR)', '1300', '14.9', '1.06', '1290', '1424')]
tbl = s.shapes.add_table(len(rows), 6, Inches(2.4), Inches(1.25),
                         Inches(8.5), Inches(4.6)).table
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        c = tbl.cell(i, j)
        c.text = val
        c.text_frame.paragraphs[0].font.size = Pt(14)
        if i == 0:
            c.text_frame.paragraphs[0].font.bold = True
add_caption(s, 'G ∝ V^7 assumed — request a 2-point HV scan to calibrate '
               'before trusting shifts >25 V.', top_in=6.2, size=13)

# ── 11. raw vs coinc per PMT ────────────────────────────────────────────────
s = add(L_ONLY)
s.shapes.title.text = 'Raw vs coincidence-selected plastic spectra (per PMT)'
add_picture(s, FIG / '10_pmt_gain/pmt_raw_vs_coinc.png', max_h_in=5.5)

# ── 12. requests ────────────────────────────────────────────────────────────
s = add(L_BODY)
s.shapes.title.text = 'Requests'
set_bullets(s.placeholders[1], [
    (0, 'READ OUT THE LIQUID SCINTILLATORS (behind the plastic on A/D): wall+LIQ '
        'coincidence requires passage through the plastic ⇒ true plastic MIP '
        'calibration. Only LIQA/LIQD cabled, no data'),
    (0, 'Plastic HV scan (2 points) → apply gain equalization'),
    (0, 'Cabling confirmations: top/bottom card assignment; which 4 of 20 wall bars '
        'unread; physical arm positions A–D; stack config per arm'),
    (0, 'DAQ follow-up on the SiPM-wall outage (bunches 643–2212)'),
])

# ── 13+. backup ─────────────────────────────────────────────────────────────
s = add(L_ONLY)
s.shapes.title.text = 'Backup: wall-only top–bottom coincidences'
add_picture(s, FIG / '06_07_geometry_mip/wall_topbot_spectra.png', max_h_in=3.9)
add_caption(s, 'Top & bottom SiPMs read the SAME 4-bar group, so any real energy '
               'deposit (through-going MIP, stopping e±, or a gamma Compton-scattering '
               'in the bars) lights both ends in tight coincidence — this removes noise, '
               'not physics backgrounds.  What it shows: SiPM response and gain are '
               'healthy on all four arms (MIP-scale bump 25–40 mV everywhere), and it '
               'provides an amplitude landmark for all 32 channels.  What it cannot '
               'show: whether particles traverse onward to the back plastic — so the '
               'A/D coincidence deficit remains unattributed (plastic response, '
               'absorption between wall and plastic, or flux composition).',
            top_in=5.3, size=13)

s = add(L_ONLY)
s.shapes.title.text = 'Backup: window scan & purity'
add_picture(s, FIG / '02_coincidence/window_scan_regions_WALC_PSSC.png', max_h_in=5.3)
s = add(L_ONLY)
s.shapes.title.text = 'Backup: rate stability (wall outage)'
add_picture(s, FIG / '01_signal_qa/rate_stability.png', max_h_in=5.3)

prs.save(OUT)
print(f'{OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)')
